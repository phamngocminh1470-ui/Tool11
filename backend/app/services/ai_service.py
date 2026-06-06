import io
import json
import logging
from typing import List, Dict, Any, Optional
import openai
import google.generativeai as genai
import pdfplumber
import docx
import pptx
from app.core.config import settings

logger = logging.getLogger("studyos.ai")

# Configure APIs if keys are provided
if settings.OPENAI_API_KEY and not settings.OPENAI_API_KEY.startswith("mock"):
    openai.api_key = settings.OPENAI_API_KEY

if settings.GEMINI_API_KEY and not settings.GEMINI_API_KEY.startswith("mock"):
    genai.configure(api_key=settings.GEMINI_API_KEY)


# =====================================================================
# 1. FILE PARSING SERVICE
# =====================================================================

def extract_text_from_bytes(file_bytes: bytes, filename: str, mime_type: str) -> str:
    """
    Extracts text from PDF, DOCX, PPTX, TXT or Image bytes.
    """
    text = ""
    try:
        if "pdf" in mime_type or filename.endswith(".pdf"):
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                text_pages = [page.extract_text() or "" for page in pdf.pages]
                text = "\n".join(text_pages)
                
        elif "word" in mime_type or filename.endswith(".docx"):
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join([p.text for p in doc.paragraphs])
            
        elif "presentation" in mime_type or filename.endswith(".pptx"):
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text_runs.append(paragraph.text)
            text = "\n".join(text_runs)
            
        elif "text" in mime_type or filename.endswith(".txt"):
            try:
                text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                text = file_bytes.decode("latin-1")
                
        elif "image" in mime_type or filename.endswith((".png", ".jpg", ".jpeg")):
            # If AI is enabled and has keys, we can send vision requests.
            # Otherwise we output a descriptive text for mock analysis.
            text = f"[Image OCR Content from {filename}]: Phân tích trực quan của biểu đồ hoặc văn bản trong ảnh."
            
        else:
            text = f"[Văn bản thô từ {filename}]: Định dạng không xác định được chuyển đổi."
            
    except Exception as e:
        logger.error(f"Error parsing file {filename}: {str(e)}")
        text = f"[Lỗi đọc tài liệu]: Không thể trích xuất văn bản tự động từ {filename}. Chi tiết lỗi: {str(e)}"
        
    return text.strip() if text else f"[Tài liệu rỗng]: Không tìm thấy nội dung trong {filename}."


# =====================================================================
# 2. AI COMPLETION ROUTER (OPENAI / GEMINI / MOCK)
# =====================================================================

def call_llm(prompt: str, system_prompt: str = "You are StudyOS, a world-class AI tutor.", response_format: str = "text") -> str:
    """
    Helper to route completion calls between OpenAI, Gemini, or Mock fallback.
    """
    # 1. Attempt Gemini
    if settings.GEMINI_API_KEY and not settings.GEMINI_API_KEY.startswith("mock"):
        try:
            model = genai.GenerativeModel(
                model_name="gemini-1.5-flash",
                system_instruction=system_prompt
            )
            response = model.generate_content(prompt)
            return response.text
        except Exception as e:
            logger.warning(f"Gemini API execution failed, trying OpenAI: {str(e)}")

    # 2. Attempt OpenAI
    if settings.OPENAI_API_KEY and not settings.OPENAI_API_KEY.startswith("mock"):
        try:
            client = openai.OpenAI(api_key=settings.OPENAI_API_KEY)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.2
            )
            return response.choices[0].message.content
        except Exception as e:
            logger.warning(f"OpenAI API execution failed: {str(e)}")

    # 3. Raise to fall back to mock generation
    raise ValueError("No valid AI API Keys configured. Falling back to dynamic Mock generation.")


# =====================================================================
# 3. DYNAMIC MOCK GENERATOR (Runs when API keys are not valid/provided)
# =====================================================================

def get_topic_context(filename: str) -> Dict[str, Any]:
    """
    Infers the subject topic based on filename to generate realistic mock data.
    """
    fn = filename.lower()
    if "marketing" in fn or "business" in fn or "sales" in fn:
        return {
            "subject": "Kinh doanh & Tiếp thị",
            "chapters": [
                {
                    "title": "Chương 1: Tổng quan về Marketing hỗn hợp (4P)",
                    "summary": "Giới thiệu mô hình 4P: Product, Price, Place, Promotion. Đây là nền tảng của mọi chiến lược tiếp thị.",
                    "key_points": ["Định nghĩa từng chữ P trong mô hình.", "Mối liên hệ giữa 4P và trải nghiệm khách hàng.", "Lịch sử hình thành mô hình."],
                    "formulas": ["ROI = (Doanh thu - Chi phí) / Chi phí"],
                    "keywords": ["Marketing Mix", "Sản phẩm", "Giá cả", "Phân phối", "Chiêu thị"]
                },
                {
                    "title": "Chương 2: Phân khúc thị trường và Định vị thương hiệu",
                    "summary": "Kỹ thuật chia nhỏ thị trường dựa trên nhân khẩu học, tâm lý học và hành vi để chọn thị trường mục tiêu phù hợp.",
                    "key_points": ["Khái niệm STP (Segmentation, Targeting, Positioning).", "Bản đồ định vị thương hiệu.", "Xác định chân dung khách hàng (Persona)."],
                    "formulas": ["LTV (Giá trị vòng đời khách hàng) = Chi tiêu trung bình * Tần suất mua * Thời gian gắn bó"],
                    "keywords": ["Phân khúc thị trường", "STP", "Khách hàng mục tiêu", "Định vị", "USP"]
                }
            ],
            "flashcards": [
                {"question": "Marketing Mix là gì?", "answer": "Tập hợp các công cụ tiếp thị mà doanh nghiệp sử dụng để đạt mục tiêu trong thị trường.", "definition": "Mô hình kết hợp các yếu tố Product, Price, Place, Promotion.", "category": "Marketing"},
                {"question": "USP là gì?", "answer": "Unique Selling Proposition - Điểm bán hàng độc nhất của sản phẩm giúp phân biệt với đối thủ.", "definition": "Đặc tính độc đáo duy nhất của sản phẩm.", "category": "Thương hiệu"},
                {"question": "STP trong marketing nghĩa là gì?", "answer": "Segmentation (Phân khúc), Targeting (Nhắm mục tiêu), Positioning (Định vị).", "definition": "Quy trình ba bước cốt lõi thiết kế chiến lược định vị.", "category": "Chiến lược"}
            ],
            "quizzes": [
                {
                    "type": "multiple_choice",
                    "difficulty": "medium",
                    "question": "Chữ P nào liên quan đến việc xác định kênh phân phối sản phẩm?",
                    "options": ["Product", "Price", "Place", "Promotion"],
                    "correct_answer": "Place",
                    "explanation": "Place đại diện cho địa điểm, các kênh phân phối đưa sản phẩm tiếp cận khách hàng mục tiêu."
                },
                {
                    "type": "true_false",
                    "difficulty": "easy",
                    "question": "Công thức tính ROI phản ánh tỷ suất hoàn vốn đầu tư.",
                    "options": ["Đúng", "Sai"],
                    "correct_answer": "Đúng",
                    "explanation": "ROI = Return on Investment đo lường lợi nhuận thu được so với chi phí bỏ ra."
                },
                {
                    "type": "fill_blank",
                    "difficulty": "hard",
                    "question": "Quy trình STP gồm Phân khúc, Nhắm mục tiêu và ______.",
                    "options": [],
                    "correct_answer": "Định vị",
                    "explanation": "Định vị (Positioning) là bước cuối cùng trong quy trình STP."
                }
            ]
        }
    elif "code" in fn or "python" in fn or "java" in fn or "web" in fn or "program" in fn or "data" in fn or "tech" in fn:
        return {
            "subject": "Khoa học Máy tính & Lập trình",
            "chapters": [
                {
                    "title": "Chương 1: Cấu trúc dữ liệu tuyến tính",
                    "summary": "Khái niệm và cách cài đặt Mảng (Array), Danh sách liên kết (Linked List), Ngăn xếp (Stack), Hàng đợi (Queue).",
                    "key_points": ["Độ phức tạp O(1) khi truy cập mảng.", "Linked List tối ưu cho thao tác thêm/xóa ở đầu.", "Stack hoạt động theo cơ chế LIFO, Queue theo FIFO."],
                    "formulas": ["Time Complexity: O(1) - Access, O(N) - Search"],
                    "keywords": ["Cấu trúc dữ liệu", "Array", "Linked List", "Stack", "Queue"]
                },
                {
                    "title": "Chương 2: Giải thuật sắp xếp và tìm kiếm",
                    "summary": "Tìm hiểu các thuật toán sắp xếp phổ biến như Bubble Sort, Quick Sort, Merge Sort và thuật toán Tìm kiếm nhị phân (Binary Search).",
                    "key_points": ["Quick Sort sử dụng chiến lược chia để trị (Divide and Conquer).", "Binary Search yêu cầu mảng phải được sắp xếp trước.", "Độ phức tạp trung bình của Quick Sort là O(N log N)."],
                    "formulas": ["Binary Search Time Complexity: O(log N)"],
                    "keywords": ["Giải thuật", "Sắp xếp", "Quick Sort", "Binary Search", "Chia để trị"]
                }
            ],
            "flashcards": [
                {"question": "Đặc điểm của Ngăn xếp (Stack) là gì?", "answer": "Hoạt động theo nguyên lý LIFO (Last In First Out) - Vào sau ra trước.", "definition": "Cấu trúc dữ liệu lưu trữ tuyến tính có giới hạn thao tác ở một đầu.", "category": "DSA"},
                {"question": "Độ phức tạp thời gian của Tìm kiếm nhị phân là bao nhiêu?", "answer": "O(log N)", "definition": "Hiệu suất thuật toán tối ưu trên tập dữ liệu đã sắp xếp.", "category": "Giải thuật"},
                {"question": "Mảng có ưu điểm gì lớn nhất?", "answer": "Truy xuất phần tử theo chỉ số (index) với thời gian O(1).", "definition": "Cấu trúc lưu trữ các phần tử cùng kiểu liên tiếp trong bộ nhớ.", "category": "DSA"}
            ],
            "quizzes": [
                {
                    "type": "multiple_choice",
                    "difficulty": "medium",
                    "question": "Cơ chế hoạt động của Hàng đợi (Queue) là gì?",
                    "options": ["LIFO", "FIFO", "LILO", "Ngẫu nhiên"],
                    "correct_answer": "FIFO",
                    "explanation": "FIFO (First In First Out) nghĩa là phần tử nào vào hàng đợi trước sẽ được lấy ra trước."
                },
                {
                    "type": "true_false",
                    "difficulty": "easy",
                    "question": "Thuật toán tìm kiếm nhị phân hoạt động tốt ngay cả trên mảng chưa được sắp xếp.",
                    "options": ["Đúng", "Sai"],
                    "correct_answer": "Sai",
                    "explanation": "Tìm kiếm nhị phân bắt buộc mảng đầu vào phải được sắp xếp theo thứ tự nhất định."
                },
                {
                    "type": "fill_blank",
                    "difficulty": "medium",
                    "question": "Thuật toán sắp xếp Quick Sort hoạt động dựa trên phương pháp ______.",
                    "options": [],
                    "correct_answer": "Chia để trị",
                    "explanation": "Quick Sort phân rã bài toán lớn thành các bài toán con xung quanh một phần tử chốt (pivot)."
                }
            ]
        }
    else:
        # Default fallback: General Education Study Guide
        return {
            "subject": "Phương pháp nghiên cứu & Học tập",
            "chapters": [
                {
                    "title": "Chương 1: Cách thiết lập mục tiêu SMART",
                    "summary": "Phương pháp đặt mục tiêu khoa học giúp tăng năng suất học tập: Specific, Measurable, Actionable, Relevant, Time-bound.",
                    "key_points": ["Mục tiêu phải rõ ràng và đo lường được.", "Gắn liền thời gian cụ thể hoàn thành.", "Liên kết với mục tiêu lâu dài."],
                    "formulas": ["SMART = Cụ thể + Đo lường + Khả thi + Thực tế + Thời hạn"],
                    "keywords": ["SMART", "Đặt mục tiêu", "Năng suất", "Học tập hiệu quả"]
                },
                {
                    "title": "Chương 2: Kỹ thuật ghi nhớ Pomodoro và Active Recall",
                    "summary": "Phương pháp học ngắt quãng 25 phút tập trung kèm 5 phút nghỉ, kết hợp việc chủ động gợi nhớ kiến thức để ghi nhớ dài hạn.",
                    "key_points": ["Active Recall kích hoạt các kết nối thần kinh mạnh hơn đọc lại thụ động.", "Lặp lại ngắt quãng (Spaced Repetition) giảm thiểu đường cong quên lãng.", "Pomodoro duy trì sự tập trung tối đa."],
                    "formulas": ["Hiệu suất = Thời gian tập trung / Tổng thời gian học"],
                    "keywords": ["Pomodoro", "Active Recall", "Spaced Repetition", "Ghi nhớ"]
                }
            ],
            "flashcards": [
                {"question": "SMART đại diện cho điều gì?", "answer": "Specific (Cụ thể), Measurable (Đo lường), Achievable (Khả thi), Relevant (Thực tế), Time-bound (Giới hạn thời gian).", "definition": "Khung tiêu chuẩn thiết kế mục tiêu hiệu quả.", "category": "Kỹ năng"},
                {"question": "Active Recall là gì?", "answer": "Phương pháp học tập chủ động truy xuất kiến thức từ não bộ thay vì đọc lại tài liệu thụ động.", "definition": "Quá trình kiểm tra trí nhớ chủ động.", "category": "Phương pháp"},
                {"question": "Kỹ thuật Pomodoro hoạt động như thế nào?", "answer": "Tập trung cao độ trong 25 phút, sau đó nghỉ ngắn 5 phút.", "definition": "Chu kỳ học tập ngắt quãng giúp bảo toàn năng lượng não bộ.", "category": "Thời gian"}
            ],
            "quizzes": [
                {
                    "type": "multiple_choice",
                    "difficulty": "easy",
                    "question": "Chữ T trong mục tiêu SMART viết tắt của từ gì?",
                    "options": ["Topic", "Time-bound", "Technology", "Task"],
                    "correct_answer": "Time-bound",
                    "explanation": "Time-bound có nghĩa là giới hạn thời gian thực hiện để tạo tính kỷ luật."
                },
                {
                    "type": "true_false",
                    "difficulty": "easy",
                    "question": "Active Recall là việc bạn đọc đi đọc lại một trang sách nhiều lần để ghi nhớ.",
                    "options": ["Đúng", "Sai"],
                    "correct_answer": "Sai",
                    "explanation": "Đọc lại là ôn tập thụ động. Active Recall yêu cầu bạn tự đặt câu hỏi và tự trả lời mà không nhìn sách."
                },
                {
                    "type": "fill_blank",
                    "difficulty": "medium",
                    "question": "Phương pháp học Pomodoro khuyến nghị chu kỳ tập trung là _____ phút.",
                    "options": [],
                    "correct_answer": "25",
                    "explanation": "25 phút là khoảng thời gian tối ưu cho một phiên tập trung Pomodoro tiêu chuẩn."
                }
            ]
        }


# =====================================================================
# 4. EXPOSED API IMPLEMENTATIONS
# =====================================================================

def analyze_document_content(text: str, filename: str) -> Dict[str, Any]:
    """
    Module 4 & 5: Automatically splits document, extracts chapters, key points, formulas, keywords,
    and returns summaries.
    """
    try:
        # Prompt to direct AI output as JSON
        prompt = f"""
        Analyze the following text from the document named "{filename}".
        You must structure your output as a valid JSON object matching this schema:
        {{
            "short_summary": "A brief 2-3 sentence overview",
            "detailed_summary": "A detailed multi-paragraph overview",
            "bullet_points": ["Key point 1", "Key point 2", "Key point 3"],
            "chapters": [
                {{
                    "title": "Chapter title",
                    "content_summary": "Summary of this chapter section",
                    "key_points": ["Point 1", "Point 2"],
                    "formulas": ["Formula 1", "Formula 2"],
                    "keywords": ["Keyword 1", "Keyword 2"],
                    "order": 1
                }}
            ]
        }}
        
        Text to analyze:
        {text[:8000]}
        """
        
        system_prompt = "You are a professional research assistant. You always respond with a clean valid JSON object following the requested schema and nothing else."
        
        response_str = call_llm(prompt, system_prompt)
        # Clean response str in case of markdown block wrappings
        if "```json" in response_str:
            response_str = response_str.split("```json")[1].split("```")[0].strip()
        elif "```" in response_str:
            response_str = response_str.split("```")[1].split("```")[0].strip()
            
        return json.loads(response_str)
        
    except Exception as e:
        logger.warning(f"AI Analysis failed, falling back to mock: {str(e)}")
        # Dynamic Mock fallback
        ctx = get_topic_context(filename)
        return {
            "short_summary": f"Tài liệu '{filename}' nói về chủ đề {ctx['subject']}. Tài liệu phân tích sâu sắc các lý thuyết nền tảng, cơ chế hoạt động thực tế, giúp người học dễ dàng nắm bắt kiến thức một cách khoa học.",
            "detailed_summary": f"Nội dung tài liệu '{filename}' tập trung khai thác các khía cạnh chuyên sâu của {ctx['subject']}. Người học sẽ được tiếp cận từ những khái niệm cơ bản nhất ở chương đầu đến các mô hình thực hành nâng cao và cách thức tối ưu hóa quy trình. Tài liệu trình bày rõ ràng, mạch lạc, kết hợp sơ đồ, công thức thực tế hỗ trợ tối đa việc ôn tập lâu dài.",
            "bullet_points": [
                "Cung cấp nền tảng kiến thức hệ thống.",
                "Trích xuất các công thức tính toán và ghi nhớ cốt lõi.",
                "Phương pháp thực hành kèm ví dụ minh họa trực quan.",
                "Hệ thống từ khóa chuyên ngành giúp tra cứu nhanh chóng."
            ],
            "chapters": ctx["chapters"]
        }


def generate_knowledge_map(chapters: List[Dict[str, Any]], filename: str) -> Dict[str, Any]:
    """
    Module 6: Generates Mindmap nodes, relations, a history timeline, and a Mermaid syntax string.
    """
    try:
        # Prompt to generate structured relations
        prompt = f"""
        Based on the chapters listed below, construct a knowledge map in JSON format containing:
        1. nodes: List of dicts with id, label, type (root, chapter, subtopic), and parent_id
        2. timeline: List of historical/logical events with event, time, and description
        3. relationships: List of connections with source, target, and relation type
        4. mermaid_code: A valid Mermaid mindmap or graph string representing this structure.
        
        Chapters:
        {json.dumps(chapters)}
        """
        
        system_prompt = "You are a diagrams expert. Return a clean JSON matching the requested structure."
        response_str = call_llm(prompt, system_prompt)
        
        if "```json" in response_str:
            response_str = response_str.split("```json")[1].split("```")[0].strip()
        return json.loads(response_str)
        
    except Exception as e:
        logger.warning(f"AI Knowledge Map generation failed, returning mock: {str(e)}")
        
        # Build mock mindmap
        nodes = [
            {"id": "root", "label": filename, "type": "root", "parent_id": None}
        ]
        relationships = []
        timeline = [
            {"event": "Giai đoạn bắt đầu", "time": "Tuần 1", "description": "Làm quen với các thuật ngữ cơ bản và cấu trúc tổng quát của tài liệu."},
            {"event": "Giai đoạn thực hành", "time": "Tuần 2-3", "description": "Giải quyết các bài tập thực hành, áp dụng công thức vào bài toán cụ thể."},
            {"event": "Giai đoạn đánh giá", "time": "Tuần 4", "description": "Làm đề thi thử, đánh giá lỗ hổng kiến thức và củng cố ôn tập."}
        ]
        
        mermaid_lines = [
            "graph TD",
            f"    root[\"{filename}\"]"
        ]
        
        for i, ch in enumerate(chapters):
            ch_id = f"ch_{i+1}"
            nodes.append({"id": ch_id, "label": ch["title"], "type": "chapter", "parent_id": "root"})
            relationships.append({"source": "root", "target": ch_id, "relation": "chứa"})
            mermaid_lines.append(f"    root --> {ch_id}[\"{ch['title']}\"]")
            
            for j, kw in enumerate(ch.get("keywords", [])[:3]):
                kw_id = f"ch_{i+1}_kw_{j+1}"
                nodes.append({"id": kw_id, "label": kw, "type": "subtopic", "parent_id": ch_id})
                relationships.append({"source": ch_id, "target": kw_id, "relation": "khái niệm"})
                mermaid_lines.append(f"    {ch_id} --> {kw_id}[\"{kw}\"]")
                
        return {
            "nodes": nodes,
            "timeline": timeline,
            "relationships": relationships,
            "mermaid_code": "\n".join(mermaid_lines)
        }


def generate_flashcards_ai(text: str, filename: str) -> List[Dict[str, Any]]:
    """
    Module 7: Automatically generates Flashcards (Question, Answer, Definition) based on text.
    """
    try:
        prompt = f"""
        Generate a list of 5-8 flashcards based on the text below.
        Return a valid JSON array of objects, where each object has:
        - question: The active recall question
        - answer: The concise answer
        - definition: Detailed definition or context
        - category: Subject category label
        
        Text:
        {text[:5000]}
        """
        system_prompt = "You are a study guide editor. Output only a valid JSON array of flashcard objects."
        response_str = call_llm(prompt, system_prompt)
        if "```json" in response_str:
            response_str = response_str.split("```json")[1].split("```")[0].strip()
        return json.loads(response_str)
    except Exception as e:
        logger.warning(f"AI Flashcard generation failed, using mock: {str(e)}")
        return get_topic_context(filename)["flashcards"]


def generate_quizzes_ai(text: str, filename: str, difficulty: str = "medium", quiz_type: str = "mixed") -> List[Dict[str, Any]]:
    """
    Module 8: Generates Quizzes based on text with multiple choice, true_false, fill_blanks, match, essay.
    """
    try:
        prompt = f"""
        Generate a list of 5 test questions based on the text below.
        Difficulty level requested: {difficulty}
        Question type requested: {quiz_type}
        
        Return a valid JSON array of objects, where each object has:
        - type: "multiple_choice", "true_false", "fill_blank", "match", or "essay"
        - difficulty: "{difficulty}"
        - question: The question text
        - options: Array of options (nullable if not multiple_choice)
        - correct_answer: The correct answer string
        - explanation: The logic behind the correct answer
        
        Text:
        {text[:5000]}
        """
        system_prompt = "You are an academic examiner. Return only a valid JSON array of quiz objects."
        response_str = call_llm(prompt, system_prompt)
        if "```json" in response_str:
            response_str = response_str.split("```json")[1].split("```")[0].strip()
        return json.loads(response_str)
    except Exception as e:
        logger.warning(f"AI Quiz generation failed, using mock: {str(e)}")
        # Filter mock quiz by difficulty if needed or return all
        return get_topic_context(filename)["quizzes"]


def ask_ai_tutor(document_text: str, question: str, chat_history: List[Dict[str, str]]) -> Dict[str, Any]:
    """
    Module 10: AI Tutor RAG Q&A with citations, source details, and context highlights.
    """
    try:
        history_str = "\n".join([f"{msg['role']}: {msg['content']}" for msg in chat_history])
        
        prompt = f"""
        You are a friendly AI tutor named StudyOS. Answer the student's question based strictly on the Document Context provided below.
        If the answer cannot be found in the context, use your intelligence but clearly state it is general knowledge, not from the document.
        
        You must return a valid JSON object matching this schema:
        {{
            "answer": "Detailed answer formatted in markdown, explaining clearly.",
            "citations": [
                {{
                    "source": "Name of section or page reference",
                    "page": 1,
                    "context": "Exact sentence or paragraph snippet quote from the document supporting the answer"
                }}
            ]
        }}
        
        Student Question:
        {question}
        
        Chat History:
        {history_str}
        
        Document Context:
        {document_text[:6000]}
        """
        
        system_prompt = "You are StudyOS, an expert RAG teaching assistant. Respond only in the requested JSON structure."
        response_str = call_llm(prompt, system_prompt)
        
        if "```json" in response_str:
            response_str = response_str.split("```json")[1].split("```")[0].strip()
        return json.loads(response_str)
        
    except Exception as e:
        logger.warning(f"AI Tutor query failed, returning structured mock: {str(e)}")
        # Simulated answering
        return {
            "answer": f"Dựa trên nội dung tài liệu của bạn, câu hỏi **'{question}'** liên quan đến các khái niệm cốt lõi. Tôi đã định vị được các từ khóa quan trọng và đề xuất bạn tập trung nghiên cứu kỹ hơn sơ đồ chương liên quan. Bạn có muốn tôi giải thích chi tiết hơn về phần này không?",
            "citations": [
                {
                    "source": "Trích lục văn bản tài liệu",
                    "page": 1,
                    "context": "Văn bản mẫu được đối chiếu từ tài liệu gốc nhằm trả lời câu hỏi học tập của bạn."
                }
            ]
        }
